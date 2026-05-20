"""
Convert documents to individual markdown files in wiki.
Handles spreadsheets by converting each visible tab to a separate markdown file.
Handles CSV files by converting them to markdown tables.
For other document types, can be extended to use DocumentConverter.

Usage: python alltomd.py <input_file> [output_dir]
  input_file: path to any document (spreadsheet, CSV, Word, PDF, etc)
  output_dir: destination folder (default: ../../wiki/)

Output naming:
  Spreadsheets: {original_filename}_{tab_name}.md for each visible tab
  CSV: {original_filename}.md
  Documents: {original_filename}.md
"""

import sys
import csv
import json
import hashlib
from datetime import datetime, timezone
from pathlib import Path
from openpyxl import load_workbook

try:
    from docx import Document as DocxDocument
except ImportError:
    DocxDocument = None

try:
    from pptx import Presentation
except ImportError:
    Presentation = None


def safe_filename(name: str) -> str:
    """Replace Windows-illegal characters with underscore, preserve spaces and hyphens."""
    illegal_chars = '<>:"|?*\\'
    result = ''.join('_' if c in illegal_chars else c for c in name)
    return result.rstrip('.')


def compute_file_hash(file_path: Path) -> str:
    """Compute SHA256 hash of a file."""
    sha256_hash = hashlib.sha256()
    with open(file_path, 'rb') as f:
        for byte_block in iter(lambda: f.read(4096), b""):
            sha256_hash.update(byte_block)
    return sha256_hash.hexdigest()


def read_manifest(output_dir: Path) -> dict:
    """Read manifest.json from wiki directory. Return empty dict if not found."""
    manifest_path = output_dir / "manifest.json"
    if manifest_path.exists():
        try:
            with open(manifest_path, 'r', encoding='utf-8') as f:
                return json.load(f)
        except (json.JSONDecodeError, IOError):
            return {}
    return {}


def write_manifest(output_dir: Path, manifest: dict) -> None:
    """Write manifest.json to wiki directory."""
    manifest_path = output_dir / "manifest.json"
    with open(manifest_path, 'w', encoding='utf-8') as f:
        json.dump(manifest, f, indent=2)


def should_convert(wiki_file: str, source_file: Path, output_dir: Path, manifest: dict) -> bool:
    """Check if a source file has changed since last conversion."""
    if wiki_file not in manifest:
        return True

    current_hash = compute_file_hash(source_file)
    return manifest[wiki_file].get('source_hash') != current_hash


def update_manifest_entry(manifest: dict, wiki_file: str, source_file: Path, source_tab: str = None) -> None:
    """Update a manifest entry with current file info."""
    manifest[wiki_file] = {
        'wiki_file': wiki_file,
        'source_file': source_file.name,
        'source_tab': source_tab,
        'converted_at': datetime.now(timezone.utc).isoformat(),
        'source_hash': compute_file_hash(source_file)
    }


def csv_to_markdown(input_path: Path, output_dir: Path, manifest: dict) -> None:
    """Convert CSV file to markdown table."""
    input_path = Path(input_path).resolve()
    output_dir = Path(output_dir).resolve()

    if not input_path.exists():
        print(f"Error: input file not found: {input_path}", file=sys.stderr)
        sys.exit(1)

    output_dir.mkdir(parents=True, exist_ok=True)

    wiki_file = f"{input_path.stem}.md"

    # Check if conversion is needed
    if not should_convert(wiki_file, input_path, output_dir, manifest):
        print(f"No change — skipping {input_path.name}")
        return

    # Read CSV file
    rows = []
    with open(input_path, 'r', encoding='utf-8') as f:
        reader = csv.reader(f)
        rows = list(reader)

    if not rows:
        print(f"Error: CSV file is empty: {input_path}", file=sys.stderr)
        return

    # Use first row as headers
    headers = [h.strip() if h else "" for h in rows[0]]
    headers = [h.replace("|", "\\|").replace("\n", " ") for h in headers]

    # Build markdown table
    lines = []
    lines.append("| " + " | ".join(headers) + " |")
    lines.append("| " + " | ".join(["---"] * len(headers)) + " |")

    # Process data rows
    for row in rows[1:]:
        cells = []
        for value in row:
            cell_str = value.strip() if value else ""
            cell_str = cell_str.replace("|", "\\|").replace("\n", " ")
            cells.append(cell_str)

        # Skip trailing all-empty rows
        if not any(cells):
            continue

        lines.append("| " + " | ".join(cells) + " |")

    markdown_content = "\n".join(lines)

    # Write to file: {csv_filename}.md
    output_path = output_dir / f"{input_path.stem}.md"

    with open(output_path, 'w', encoding='utf-8') as f:
        f.write(markdown_content)
        f.write("\n")

    update_manifest_entry(manifest, wiki_file, input_path)
    print(f"Converted: {input_path.name} -> {output_path.name}")


def sheet_to_markdown(ws) -> str:
    """Convert openpyxl worksheet to markdown table string."""
    rows = list(ws.iter_rows(values_only=True))

    if not rows:
        return ""

    # Use first row as headers
    headers = [str(h) if h is not None else "" for h in rows[0]]

    # Escape pipes and collapse newlines in headers
    headers = [h.replace("|", "\\|").replace("\n", " ").strip() for h in headers]

    # Build markdown table
    lines = []
    lines.append("| " + " | ".join(headers) + " |")
    lines.append("| " + " | ".join(["---"] * len(headers)) + " |")

    # Process data rows
    for row in rows[1:]:
        cells = []
        for value in row:
            if value is None:
                cell_str = ""
            else:
                cell_str = str(value).replace("|", "\\|").replace("\n", " ").strip()
            cells.append(cell_str)

        # Skip trailing all-empty rows
        if not any(cells):
            continue

        lines.append("| " + " | ".join(cells) + " |")

    return "\n".join(lines)


def convert_excel_to_md(input_path: Path, output_dir: Path, manifest: dict) -> None:
    """Convert each visible sheet in Excel file to individual markdown file."""
    input_path = Path(input_path).resolve()
    output_dir = Path(output_dir).resolve()

    if not input_path.exists():
        print(f"Error: input file not found: {input_path}", file=sys.stderr)
        sys.exit(1)

    output_dir.mkdir(parents=True, exist_ok=True)

    # Check if source file has changed (for all tabs)
    file_changed = should_convert(f"{input_path.stem}_check", input_path, output_dir, {"check": manifest.get(f"{input_path.stem}_check", {})})

    if not file_changed:
        print(f"No change — skipping {input_path.name}")
        return

    wb = load_workbook(input_path, data_only=True)

    converted_count = 0
    skipped_count = 0
    original_stem = input_path.stem

    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]

        # Skip hidden sheets
        if ws.sheet_state != "visible":
            skipped_count += 1
            print(f"Skipping hidden sheet: {sheet_name}")
            continue

        # Convert to markdown
        markdown_content = sheet_to_markdown(ws)

        # Write to file: {original_filename}_{tab_name}.md
        safe_sheet_name = safe_filename(sheet_name)
        wiki_file = f"{original_stem}_{safe_sheet_name}.md"
        output_path = output_dir / wiki_file

        with open(output_path, 'w', encoding='utf-8') as f:
            f.write(f"# {sheet_name}\n\n")
            f.write(markdown_content)
            f.write("\n")

        update_manifest_entry(manifest, wiki_file, input_path, source_tab=sheet_name)
        print(f"Converted: {sheet_name} -> {output_path.name}")
        converted_count += 1

    print(f"\nSummary: {converted_count} sheets converted, {skipped_count} hidden sheets skipped.")


def _paragraph_to_markdown(para) -> str:
    """Convert a python-docx paragraph to markdown string."""
    if not para.text.strip():
        return ""

    # Map style names to heading levels
    style_name = para.style.name if para.style else ""

    if style_name.startswith("Heading"):
        level = 1
        if "Heading 1" in style_name:
            level = 1
        elif "Heading 2" in style_name:
            level = 2
        elif "Heading 3" in style_name:
            level = 3
        elif "Heading 4" in style_name:
            level = 4
        elif "Heading 5" in style_name:
            level = 5
        elif "Heading 6" in style_name:
            level = 6
        return "#" * level + " " + para.text

    # Handle list styles
    if "List" in style_name:
        indent = int(para.paragraph_format.left_indent.pt // 36) if para.paragraph_format.left_indent else 0
        return "  " * indent + "- " + para.text

    # Regular paragraph: apply inline formatting
    result = ""
    for run in para.runs:
        text = run.text
        if run.bold:
            text = f"**{text}**"
        if run.italic:
            text = f"*{text}*"
        result += text

    return result if result.strip() else ""


def _docx_table_to_markdown(table) -> str:
    """Convert a python-docx table to markdown table string."""
    lines = []

    for i, row in enumerate(table.rows):
        cells = []
        for cell in row.cells:
            cell_text = cell.text.replace("|", "\\|").replace("\n", " ").strip()
            cells.append(cell_text)

        lines.append("| " + " | ".join(cells) + " |")

        # Add separator after header row (first row)
        if i == 0:
            lines.append("| " + " | ".join(["---"] * len(cells)) + " |")

    return "\n".join(lines)


def docx_to_markdown(input_path: Path, output_dir: Path, manifest: dict) -> None:
    """Convert Word document to markdown file."""
    if DocxDocument is None:
        print("Error: python-docx not installed. Run: pip install python-docx", file=sys.stderr)
        sys.exit(1)

    input_path = Path(input_path).resolve()
    output_dir = Path(output_dir).resolve()

    if not input_path.exists():
        print(f"Error: input file not found: {input_path}", file=sys.stderr)
        sys.exit(1)

    output_dir.mkdir(parents=True, exist_ok=True)

    wiki_file = f"{input_path.stem}.md"

    # Check if conversion is needed
    if not should_convert(wiki_file, input_path, output_dir, manifest):
        print(f"No change — skipping {input_path.name}")
        return

    # Open document
    doc = DocxDocument(input_path)

    # Build maps for element lookup
    from docx.oxml.text.paragraph import CT_P
    from docx.table import _Cell

    para_map = {p._element: p for p in doc.paragraphs}
    table_map = {t._element: t for t in doc.tables}

    # Convert content
    lines = [f"# {input_path.stem}\n"]

    for element in doc.element.body:
        if element in para_map:
            md_text = _paragraph_to_markdown(para_map[element])
            if md_text.strip():
                lines.append(md_text)
        elif element in table_map:
            lines.append(_docx_table_to_markdown(table_map[element]))

    markdown_content = "\n".join(lines)

    # Write to file
    output_path = output_dir / f"{input_path.stem}.md"
    with open(output_path, 'w', encoding='utf-8') as f:
        f.write(markdown_content)
        f.write("\n")

    update_manifest_entry(manifest, wiki_file, input_path)
    print(f"Converted: {input_path.name} -> {output_path.name}")


def pptx_to_markdown(input_path: Path, output_dir: Path, manifest: dict) -> None:
    """Convert PowerPoint presentation to markdown file."""
    if Presentation is None:
        print("Error: python-pptx not installed. Run: pip install python-pptx", file=sys.stderr)
        sys.exit(1)

    input_path = Path(input_path).resolve()
    output_dir = Path(output_dir).resolve()

    if not input_path.exists():
        print(f"Error: input file not found: {input_path}", file=sys.stderr)
        sys.exit(1)

    output_dir.mkdir(parents=True, exist_ok=True)

    wiki_file = f"{input_path.stem}.md"

    # Check if conversion is needed
    if not should_convert(wiki_file, input_path, output_dir, manifest):
        print(f"No change — skipping {input_path.name}")
        return

    # Open presentation
    prs = Presentation(input_path)

    lines = [f"# {input_path.stem}\n"]

    for slide_num, slide in enumerate(prs.slides, 1):
        lines.append(f"## Slide {slide_num}\n")

        for shape in slide.shapes:
            # Extract text from shape
            if hasattr(shape, "text") and shape.text.strip():
                lines.append(shape.text)

            # Extract table from shape
            if shape.has_table:
                table = shape.table
                lines.append(_docx_table_to_markdown(table))

    markdown_content = "\n".join(lines)

    # Write to file
    output_path = output_dir / f"{input_path.stem}.md"
    with open(output_path, 'w', encoding='utf-8') as f:
        f.write(markdown_content)
        f.write("\n")

    update_manifest_entry(manifest, wiki_file, input_path)
    print(f"Converted: {input_path.name} -> {output_path.name}")


def main():
    """Main entry point. Input path is required."""
    script_dir = Path(__file__).parent
    default_output = script_dir.parent.parent.parent / "wiki"

    if len(sys.argv) < 2:
        print("Usage: python alltomd.py <input_file> [output_dir]", file=sys.stderr)
        print(f"Default output_dir: {default_output}", file=sys.stderr)
        sys.exit(1)

    input_path = Path(sys.argv[1])
    output_dir = Path(sys.argv[2] if len(sys.argv) > 2 else default_output)

    # Load existing manifest
    output_dir.mkdir(parents=True, exist_ok=True)
    manifest = read_manifest(output_dir)

    # Detect file type and call appropriate converter
    if input_path.suffix.lower() == '.csv':
        csv_to_markdown(input_path, output_dir, manifest)
    elif input_path.suffix.lower() in ['.xlsx', '.xls', '.xlsm']:
        convert_excel_to_md(input_path, output_dir, manifest)
    elif input_path.suffix.lower() == '.docx':
        docx_to_markdown(input_path, output_dir, manifest)
    elif input_path.suffix.lower() == '.pptx':
        pptx_to_markdown(input_path, output_dir, manifest)
    else:
        print(f"Error: unsupported file type: {input_path.suffix}", file=sys.stderr)
        print("Supported: .csv, .xlsx, .xlsm, .xls, .docx, .pptx", file=sys.stderr)
        sys.exit(1)

    # Write manifest back
    write_manifest(output_dir, manifest)


if __name__ == "__main__":
    main()
