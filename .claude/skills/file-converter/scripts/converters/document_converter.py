"""
Document to Markdown Converter
Handles .docx, .pptx, .pdf, .txt files
"""

import os
import base64
from pathlib import Path
from typing import List, Tuple
import re

try:
    from docx import Document as DocxDocument
    from docx.oxml.text.paragraph import CT_P
    from docx.table import _Cell, Table
except ImportError:
    DocxDocument = None

try:
    from pptx import Presentation
except ImportError:
    Presentation = None

try:
    from PyPDF2 import PdfReader
except ImportError:
    PdfReader = None


class DocumentConverter:
    """Converts various document formats to Markdown."""

    def __init__(self, output_dir: Path = None):
        self.output_dir = Path(output_dir) if output_dir else Path('/mnt/user-data/outputs')
        self.output_dir.mkdir(parents=True, exist_ok=True)
    
    def convert(self, input_path: Path) -> Path:
        """Convert document to markdown."""
        file_ext = input_path.suffix.lower()
        
        if file_ext == '.docx':
            content = self._convert_docx(input_path)
        elif file_ext == '.pptx':
            content = self._convert_pptx(input_path)
        elif file_ext == '.pdf':
            content = self._convert_pdf(input_path)
        elif file_ext == '.txt':
            content = self._convert_txt(input_path)
        else:
            raise ValueError(f"Unsupported document format: {file_ext}")
        
        # Write output
        output_path = self.output_dir / f"{input_path.stem}.md"
        output_path.write_text(content, encoding='utf-8')
        return output_path
    
    def _convert_docx(self, input_path: Path) -> str:
        """Convert DOCX to Markdown."""
        if DocxDocument is None:
            raise ImportError("python-docx is required. Install with: pip install python-docx")
        
        doc = DocxDocument(input_path)
        md_content = []
        
        for element in doc.element.body:
            if isinstance(element, CT_P):
                # Handle paragraph
                para = element.getparent().find('.//{http://schemas.openxmlformats.org/wordprocessingml/2006/main}p')
                for para_elem in doc.paragraphs:
                    if para_elem._element == element:
                        md_content.append(self._paragraph_to_markdown(para_elem))
                        break
            elif element.tag.endswith('tbl'):
                # Handle table
                for table in doc.tables:
                    if table._element == element:
                        md_content.append(self._table_to_markdown(table))
                        break
        
        # Fallback: iterate through paragraphs and tables
        if not md_content:
            for para in doc.paragraphs:
                md_content.append(self._paragraph_to_markdown(para))
            
            for table in doc.tables:
                md_content.append(self._table_to_markdown(table))
        
        return '\n\n'.join(filter(None, md_content))
    
    def _paragraph_to_markdown(self, para) -> str:
        """Convert a paragraph to markdown."""
        text = ''
        style = para.style.name if para.style else 'Normal'
        
        # Get text with formatting
        for run in para.runs:
            run_text = run.text
            if run.bold:
                run_text = f"**{run_text}**"
            if run.italic:
                run_text = f"*{run_text}*"
            text += run_text
        
        if not text.strip():
            return ''
        
        # Apply paragraph styles
        if 'Heading 1' in style:
            return f"# {text}"
        elif 'Heading 2' in style:
            return f"## {text}"
        elif 'Heading 3' in style:
            return f"### {text}"
        elif 'Heading 4' in style:
            return f"#### {text}"
        elif 'List' in style or para.style.name.startswith('List'):
            level = para.paragraph_format.left_indent.pt // 36 if para.paragraph_format.left_indent else 0
            indent = '  ' * level
            return f"{indent}- {text}"
        else:
            return text
    
    def _table_to_markdown(self, table) -> str:
        """Convert a table to markdown."""
        md_lines = []
        
        for i, row in enumerate(table.rows):
            cells = []
            for cell in row.cells:
                cell_text = ''.join(p.text for p in cell.paragraphs).strip()
                cells.append(cell_text)
            
            md_lines.append('| ' + ' | '.join(cells) + ' |')
            
            # Add separator after header row
            if i == 0:
                separators = ['---'] * len(cells)
                md_lines.append('| ' + ' | '.join(separators) + ' |')
        
        return '\n'.join(md_lines)
    
    def _convert_pptx(self, input_path: Path) -> str:
        """Convert PPTX to Markdown."""
        if Presentation is None:
            raise ImportError("python-pptx is required. Install with: pip install python-pptx")
        
        prs = Presentation(input_path)
        md_content = []
        
        for slide_num, slide in enumerate(prs.slides, 1):
            md_content.append(f"## Slide {slide_num}")
            
            # Extract text from shapes
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    md_content.append(shape.text)
                
                # Extract table content
                if shape.has_table:
                    md_content.append(self._pptx_table_to_markdown(shape.table))
                
                # Extract notes
                if hasattr(shape, "notes_slide"):
                    notes = shape.notes_slide.notes_text_frame.text.strip()
                    if notes:
                        md_content.append(f"*Note: {notes}*")
            
            md_content.append('')
        
        return '\n\n'.join(filter(None, md_content))
    
    def _pptx_table_to_markdown(self, table) -> str:
        """Convert PPTX table to markdown."""
        md_lines = []
        
        for i, row in enumerate(table.rows):
            cells = []
            for cell in row.cells:
                cell_text = cell.text.strip()
                cells.append(cell_text)
            
            md_lines.append('| ' + ' | '.join(cells) + ' |')
            
            if i == 0:
                separators = ['---'] * len(cells)
                md_lines.append('| ' + ' | '.join(separators) + ' |')
        
        return '\n'.join(md_lines)
    
    def _convert_pdf(self, input_path: Path) -> str:
        """Convert PDF to Markdown."""
        if PdfReader is None:
            raise ImportError("PyPDF2 is required. Install with: pip install PyPDF2")
        
        reader = PdfReader(input_path)
        md_content = []
        
        for page_num, page in enumerate(reader.pages, 1):
            text = page.extract_text()
            if text.strip():
                md_content.append(f"## Page {page_num}\n")
                md_content.append(text)
        
        return '\n\n'.join(filter(None, md_content))
    
    def _convert_txt(self, input_path: Path) -> str:
        """Convert TXT to Markdown (minimal conversion)."""
        content = input_path.read_text(encoding='utf-8', errors='replace')
        # Basic markdown conversion for text files
        return content
