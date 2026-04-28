"""
Deterministic document preprocessor for the doc-ingestion agent.

Reads heterogeneous documents (xlsx, docx, pdf, pptx, csv) and writes a canonical
JSON shape to converted_documents/<stem>.preprocessed.json. Strips noise
(empty rows/cols, repeated headers/footers, page numbers, disclaimer blocks)
so the LLM only sees content that matters.

CLI:
    python preprocess.py <path> [<path>...]
    python preprocess.py --force <path>      # bypass mtime cache

Output JSON shape:
{
  "file": "<absolute source path>",
  "type": "xlsx|docx|pdf|pptx|csv",
  "ingested_at": "<ISO 8601>",
  "sections": [{"name": "...", "content": "..."}],
  "tables":   [{"name": "...", "rows": [[...], ...]}]
}
"""

from __future__ import annotations

import csv
import json
import re
import sys
from datetime import datetime, timezone
from pathlib import Path

SCRIPT_DIR = Path(__file__).resolve().parent
BASE_DIR = SCRIPT_DIR.parents[3]
OUTPUT_DIR = BASE_DIR / "converted_documents"

DISCLAIMER_PHRASES = [
    "this document is confidential",
    "past performance is not indicative",
    "past performance is not a reliable indicator",
    "this information is general in nature",
    "this is not financial advice",
    "disclaimer",
    "important information",
    "all rights reserved",
    "no part of this document may be reproduced",
    "the information contained in this document",
]

PAGE_NUMBER_PATTERNS = [
    re.compile(r"^\s*page\s+\d+(\s+of\s+\d+)?\s*$", re.IGNORECASE),
    re.compile(r"^\s*\d+\s*$"),
    re.compile(r"^\s*-\s*\d+\s*-\s*$"),
]


def _is_page_number(line: str) -> bool:
    return any(p.match(line) for p in PAGE_NUMBER_PATTERNS)


def _is_disclaimer(line: str) -> bool:
    low = line.strip().lower()
    return any(phrase in low for phrase in DISCLAIMER_PHRASES)


def _strip_repeated_lines(pages: list[list[str]], threshold: float = 0.5) -> list[list[str]]:
    if len(pages) < 2:
        return pages
    counts: dict[str, int] = {}
    for page in pages:
        for ln in {l.strip() for l in page if l.strip()}:
            counts[ln] = counts.get(ln, 0) + 1
    cutoff = max(2, int(len(pages) * threshold))
    repeats = {ln for ln, c in counts.items() if c >= cutoff and len(ln) < 200}
    return [[l for l in page if l.strip() not in repeats] for page in pages]


def _clean_table_rows(rows: list[list]) -> list[list]:
    cleaned = [[("" if c is None else str(c)).strip() for c in r] for r in rows]
    cleaned = [r for r in cleaned if any(c for c in r)]
    if not cleaned:
        return []
    width = max(len(r) for r in cleaned)
    cleaned = [r + [""] * (width - len(r)) for r in cleaned]
    keep_cols = [i for i in range(width) if any(r[i] for r in cleaned)]
    return [[r[i] for i in keep_cols] for r in cleaned]


def _drop_title_rows(rows: list[list]) -> list[list]:
    """Drop leading single-cell-populated rows (typical xlsx title banners)."""
    out = list(rows)
    while out and sum(1 for c in out[0] if c) <= 1 and len(out[0]) > 1:
        out.pop(0)
    return out


def preprocess_xlsx(path: Path) -> dict:
    from openpyxl import load_workbook

    wb = load_workbook(filename=str(path), data_only=True, read_only=True)
    tables = []
    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        rows = [list(r) for r in ws.iter_rows(values_only=True)]
        rows = _clean_table_rows(rows)
        rows = _drop_title_rows(rows)
        if rows:
            tables.append({"name": sheet_name, "rows": rows})
    wb.close()
    return {"sections": [], "tables": tables}


def preprocess_docx(path: Path) -> dict:
    from docx import Document

    doc = Document(str(path))
    paragraphs = []
    for p in doc.paragraphs:
        text = p.text.strip()
        if not text or _is_page_number(text) or _is_disclaimer(text):
            continue
        paragraphs.append(text)
    sections = [{"name": "body", "content": "\n".join(paragraphs)}] if paragraphs else []
    tables = []
    for i, t in enumerate(doc.tables):
        rows = [[cell.text.strip() for cell in row.cells] for row in t.rows]
        rows = _clean_table_rows(rows)
        if rows:
            tables.append({"name": f"table_{i + 1}", "rows": rows})
    return {"sections": sections, "tables": tables}


def preprocess_pdf(path: Path) -> dict:
    try:
        import pdfplumber
    except ImportError:
        from pypdf import PdfReader

        reader = PdfReader(str(path))
        pages = [(p.extract_text() or "").splitlines() for p in reader.pages]
    else:
        with pdfplumber.open(str(path)) as pdf:
            pages = [(p.extract_text() or "").splitlines() for p in pdf.pages]

    pages = _strip_repeated_lines(pages)
    cleaned_pages = []
    for page in pages:
        kept = [l for l in page if l.strip() and not _is_page_number(l) and not _is_disclaimer(l)]
        if kept:
            cleaned_pages.append("\n".join(kept))

    sections = [{"name": f"page_{i + 1}", "content": text} for i, text in enumerate(cleaned_pages)]
    return {"sections": sections, "tables": []}


def preprocess_pptx(path: Path) -> dict:
    from pptx import Presentation

    prs = Presentation(str(path))
    sections = []
    for i, slide in enumerate(prs.slides):
        chunks = []
        for shape in slide.shapes:
            if not getattr(shape, "has_text_frame", False):
                continue
            for para in shape.text_frame.paragraphs:
                line = "".join(run.text for run in para.runs).strip()
                if line and not _is_page_number(line) and not _is_disclaimer(line):
                    chunks.append(line)
        if chunks:
            sections.append({"name": f"slide_{i + 1}", "content": "\n".join(chunks)})
    return {"sections": sections, "tables": []}


def preprocess_csv(path: Path) -> dict:
    with path.open(newline="", encoding="utf-8-sig", errors="replace") as f:
        rows = list(csv.reader(f))
    rows = _clean_table_rows(rows)
    rows = _drop_title_rows(rows)
    return {"sections": [], "tables": [{"name": path.stem, "rows": rows}] if rows else []}


DISPATCH = {
    ".xlsx": ("xlsx", preprocess_xlsx),
    ".xlsm": ("xlsx", preprocess_xlsx),
    ".docx": ("docx", preprocess_docx),
    ".pdf":  ("pdf",  preprocess_pdf),
    ".pptx": ("pptx", preprocess_pptx),
    ".csv":  ("csv",  preprocess_csv),
}


def preprocess(path: Path, force: bool = False) -> Path:
    suffix = path.suffix.lower()
    if suffix not in DISPATCH:
        raise ValueError(f"Unsupported file type: {suffix} ({path})")
    type_label, fn = DISPATCH[suffix]

    OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
    out_path = OUTPUT_DIR / f"{path.stem}.preprocessed.json"

    if not force and out_path.exists() and out_path.stat().st_mtime >= path.stat().st_mtime:
        return out_path

    body = fn(path)
    payload = {
        "file": str(path.resolve()),
        "type": type_label,
        "ingested_at": datetime.now(timezone.utc).isoformat(timespec="seconds"),
        "sections": body["sections"],
        "tables": body["tables"],
    }
    out_path.write_text(json.dumps(payload, ensure_ascii=False, indent=2), encoding="utf-8")
    return out_path


def main(argv: list[str]) -> int:
    args = list(argv)
    force = False
    if "--force" in args:
        force = True
        args.remove("--force")
    if not args:
        print("usage: preprocess.py [--force] <path> [<path>...]", file=sys.stderr)
        return 2
    for a in args:
        p = Path(a).resolve()
        if not p.exists():
            print(f"NOT FOUND: {p}", file=sys.stderr)
            return 1
        try:
            out = preprocess(p, force=force)
            print(out)
        except Exception as e:
            print(f"FAILED {p}: {e.__class__.__name__}: {e}", file=sys.stderr)
            return 1
    return 0


if __name__ == "__main__":
    sys.exit(main(sys.argv[1:]))
